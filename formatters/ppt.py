"""
PowerPoint 格式化器
将树形结构转换为 PowerPoint 演示文稿
"""

from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from .base import BaseFormatter


class PPTFormatter(BaseFormatter):
    """PowerPoint 格式化器"""

    # 配色方案
    COLORS = {
        "primary": RGBColor(0, 51, 102),      # 深蓝色
        "secondary": RGBColor(21, 101, 192),  # 明蓝色
        "accent": RGBColor(255, 152, 0),      # 橙色
        "text": RGBColor(33, 33, 33),         # 深灰色
        "light": RGBColor(245, 245, 245)      # 浅灰色
    }

    @property
    def format_name(self) -> str:
        return "ppt"

    @property
    def file_extension(self) -> str:
        return ".pptx"

    @property
    def description(self) -> str:
        return "PowerPoint 演示文稿格式"

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """添加标题幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片布局

        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.COLORS["primary"]

    def _add_content_slide(self, prs: Presentation, title: str,
                          content_items: List[str], level: int = 0) -> None:
        """添加内容幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局

        # 设置标题
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.COLORS["primary"]

        # 添加内容
        if content_items:
            # 获取内容占位符
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            for item in content_items:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = min(level, 4)  # PowerPoint 最多支持 5 级缩进
                p.font.size = Pt(18)
                p.font.color.rgb = self.COLORS["text"]

    def _build_slides_from_tree(self, prs: Presentation, nodes: List[Dict[str, Any]],
                                parent_title: str = None, level: int = 0) -> None:
        """
        递归构建幻灯片

        策略：
        - 每个一级节点创建一个幻灯片
        - 子节点作为项目符号列表
        """
        for node in nodes:
            node_text = node["text"]
            children = node.get("children", [])

            if level == 0:
                # 一级节点：创建新幻灯片
                if children:
                    # 收集所有子节点文本
                    content_items = self._collect_content_items(children)
                    self._add_content_slide(prs, node_text, content_items, level=0)
                else:
                    # 没有子节点，创建简单幻灯片
                    self._add_content_slide(prs, node_text, [], level=0)
            else:
                # 更深层级的节点在父幻灯片中作为列表项处理
                pass

    def _collect_content_items(self, nodes: List[Dict[str, Any]],
                               level: int = 0) -> List[str]:
        """收集内容项（扁平化树结构为列表）"""
        items = []
        indent = "  " * level

        for node in nodes:
            items.append(f"{indent}{node['text']}")
            if node.get("children"):
                # 递归收集子节点
                child_items = self._collect_content_items(node["children"], level + 1)
                items.extend(child_items)

        return items

    def format(self, title: str, tree_nodes: List[Dict[str, Any]],
               theme: str = "default", **options) -> Presentation:
        """
        将树形结构格式化为 PowerPoint 演示文稿

        参数:
            title: 演示文稿标题
            tree_nodes: 树形结构节点列表
            theme: 主题样式（默认：default）
            **options: 其他选项

        返回:
            Presentation 对象
        """
        # 创建演示文稿
        prs = Presentation()

        # 设置幻灯片尺寸为 16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # 添加标题幻灯片
        self._add_title_slide(prs, title)

        # 构建内容幻灯片
        self._build_slides_from_tree(prs, tree_nodes)

        return prs

    def save(self, content: Presentation, filename: str) -> None:
        """
        保存为 .pptx 文件

        参数:
            content: format() 方法返回的 Presentation 对象
            filename: 输出文件路径
        """
        content.save(filename)
        print(f"[OK] PowerPoint 文件生成成功: {filename}")
